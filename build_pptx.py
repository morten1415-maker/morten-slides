"""
Morten DS — PowerPoint-generator
=================================
On-brand .pptx der følger præsentations-brandguiden (PRESENTATIONS.md):
bærbar brutalisme på 16:9 — sorte kanter, hårde offset-shadows (ingen blur),
IBM Plex (mono = system, sans = indhold), sparsom rød/gul.

Kræver:  pip install python-pptx
Brug:    se deck.example.py — opret en Deck, kald .title()/.content()/... , kald .save()
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy
import os

_ASSETS = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets")
OES_LOGO_GREEN = os.path.join(_ASSETS, "oes-logo-green.png")   # til lyse slides
OES_LOGO_WHITE = os.path.join(_ASSETS, "oes-logo-white.png")   # til grønne sektioner

# ---------------------------------------------------------------- TOKENS
INK        = RGBColor(0x00, 0x00, 0x00)
INK_SOFT   = RGBColor(0x33, 0x33, 0x33)
MUTED      = RGBColor(0x6B, 0x6B, 0x6B)
CANVAS     = RGBColor(0xF4, 0xF4, 0xF2)
PAPER      = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT     = RGBColor(0xFE, 0xF3, 0xC7)   # gul
RED        = RGBColor(0xFF, 0x57, 0x57)
SUCCESS    = RGBColor(0x16, 0xA3, 0x4A)
LINE_SOFT  = RGBColor(0xE8, 0xE8, 0xE3)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

FONT_SANS  = "IBM Plex Sans"
FONT_MONO  = "IBM Plex Mono"

# 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN  = Inches(0.66)


# ---------------------------------------------------------------- LOW-LEVEL HELPERS
def _solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _no_fill(shape):
    shape.fill.background()


def _border(shape, color=INK, width_pt=2.0):
    ln = shape.line
    ln.color.rgb = color
    ln.width = Pt(width_pt)


def _no_border(shape):
    shape.line.fill.background()


def _offset_shadow(shape, color=INK, dist_pt=5):
    """Hård offset-shadow (ingen blur) injiceret via DrawingML XML.

    PowerPoint kan ikke lave 0-blur offset-shadow gennem python-pptx' API,
    så vi tilføjer <a:outerShdw> direkte. dir 2700000 = 45° (ned-højre)."""
    spPr = shape._element.spPr
    # Fjern evt. eksisterende effektliste
    for tag in ("a:effectLst",):
        existing = spPr.find(qn(tag))
        if existing is not None:
            spPr.remove(existing)
    effLst = spPr.makeelement(qn("a:effectLst"), {})
    shdw = effLst.makeelement(qn("a:outerShdw"), {
        "blurRad": "0",
        "dist": str(int(Pt(dist_pt))),
        "dir": "2700000",        # 45 grader = ned mod højre
        "rotWithShape": "0",
    })
    clr = shdw.makeelement(qn("a:srgbClr"), {"val": "%02X%02X%02X" % (color[0], color[1], color[2])})
    shdw.append(clr)
    effLst.append(shdw)
    spPr.append(effLst)


def _txt(tf, text, *, font=FONT_SANS, size=22, bold=False, color=INK_SOFT,
         align=PP_ALIGN.LEFT, mono_label=False, line_pct=None):
    """Sæt tekst på en text_frame's første paragraph."""
    p = tf.paragraphs[0]
    p.alignment = align
    if mono_label:
        text = text.upper()
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = FONT_MONO if mono_label else font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    if mono_label:
        _letter_spacing(run, 1.2)   # ~0.1em ved 11-16pt
    if line_pct:
        _line_spacing(p, line_pct)
    return p


def _letter_spacing(run, pt):
    """Letter-spacing i punkter (DrawingML 'spc' i 1/100 pt)."""
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(int(pt * 100)))


def _line_spacing(paragraph, pct):
    paragraph.line_spacing = pct


def _fill(shape, blocks, anchor=MSO_ANCHOR.TOP,
          ml=0.55, mr=0.55, mt=0.5, mb=0.5):
    """Læg flere tekst-afsnit ind i en forms egen text_frame med marginer.

    Tekst flyder naturligt og kan hverken overlappe eller ramme kanten.
    blocks: liste af dict {text, size, bold, color, font, mono, align, line,
                           before(pt), after(pt), spc}."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(ml); tf.margin_right = Inches(mr)
    tf.margin_top = Inches(mt); tf.margin_bottom = Inches(mb)
    for i, b in enumerate(blocks):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = b.get("align", PP_ALIGN.LEFT)
        if b.get("line"): p.line_spacing = b["line"]
        if b.get("before") is not None: p.space_before = Pt(b["before"])
        if b.get("after") is not None: p.space_after = Pt(b["after"])
        run = p.add_run()
        mono = b.get("mono", False)
        run.text = b["text"].upper() if mono else b["text"]
        f = run.font
        f.name = FONT_MONO if mono else b.get("font", FONT_SANS)
        f.size = Pt(b["size"]); f.bold = b.get("bold", False)
        f.color.rgb = b.get("color", INK_SOFT)
        if mono:
            _letter_spacing(run, b.get("spc", 1.2))


# ---------------------------------------------------------------- TEMAER
# Økonomistyrelsen-farver (fra brandguiden)
OES_GREEN = RGBColor(0x06, 0x6B, 0x43)   # identitetsfarve
OES_LAKS  = RGBColor(0xED, 0x5E, 0x66)   # supplerende koral (≈ vores signal-rød)

# Et tema styrer kun de farver der adskiller varianterne. Strukturen
# (sorte kanter, hårde shadows, canvas-bg, gul highlight) er ens.
THEMES = {
    # Primær præsentation stil
    "default": dict(
        signal=RED, success=SUCCESS, brand=INK,
        kicker_color=MUTED, cta_bg=INK,
        section_default="dark", brand_bar=False, logo_text=None,
    ),
    # Økonomistyrelsen — primær stil + grøn anker + officielt logo.
    # Grøn på: titel-accent, section-dividers, logo, titel-topkant, positive tal, CTA.
    # Kickers er muted grå (som primær) -> kun ÉN overskrift pr. slide.
    # Logo nederst til højre på alle slides med plads.
    "oes": dict(
        signal=OES_LAKS, success=OES_GREEN, brand=OES_GREEN,
        kicker_color=MUTED, cta_bg=OES_GREEN,
        section_default="brand", brand_bar=False, logo_text="Økonomistyrelsen",
        logo_light=OES_LOGO_GREEN, logo_dark=OES_LOGO_WHITE,
        show_eyebrow=False,   # kun ÉN overskrift: den sorte titel (ingen kicker-label)
    ),
}
THEMES["primær"] = THEMES["default"]
THEMES["primary"] = THEMES["default"]


# ---------------------------------------------------------------- DECK
class Deck:
    def __init__(self, deck_title="Morten DS", footer=None, theme="default",
                 logo_path=None):
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self.blank = self.prs.slide_layouts[6]   # tom layout
        self.deck_title = deck_title
        self.footer = footer or deck_title
        self._n = 0
        # Tema-farver som instans-attributter
        t = THEMES[theme] if isinstance(theme, str) else theme
        self.signal = t["signal"]
        self.success = t["success"]
        self.brand = t["brand"]
        self.kicker_color = t["kicker_color"]
        self.cta_bg = t["cta_bg"]
        self.section_default = t["section_default"]
        self.brand_bar = t["brand_bar"]
        self.logo_text = t.get("logo_text")
        self.logo_path = logo_path                # ægte logo-fil (override, valgfri)
        self.logo_light = t.get("logo_light")     # logo til lyse baggrunde
        self.logo_dark = t.get("logo_dark")       # logo til mørke/grønne baggrunde
        self.brand_accent = self.brand != INK     # bruger temaet en brand-farve?
        self.show_eyebrow = t.get("show_eyebrow", True)   # kicker-label over titlen?

    # --- intern: nyt blankt slide med valgfri baggrund ---
    def _new(self, bg=CANVAS):
        self._n += 1
        slide = self.prs.slides.add_slide(self.blank)
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
        _solid(rect, bg)
        _no_border(rect)
        rect.shadow.inherit = False
        # Brand-bar: tynd ØS-grøn stribe i venstre kant (kun tema med brand_bar)
        if self.brand_bar:
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.14), SLIDE_H)
            _solid(bar, self.brand)
            _no_border(bar)
            bar.shadow.inherit = False
        # 3px sort yderramme
        frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
        _no_fill(frame)
        _border(frame, INK, 3.0)
        frame.shadow.inherit = False
        return slide

    def _box(self, slide, l, t, w, h):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tb.text_frame.word_wrap = True
        return tb

    def _logo_path_for(self, dark):
        """Vælg logo-fil: brugerens override, ellers tema-logo (hvid på mørk/grøn)."""
        if self.logo_path:
            return self.logo_path
        if dark and self.logo_dark and os.path.exists(self.logo_dark):
            return self.logo_dark
        if (not dark) and self.logo_light and os.path.exists(self.logo_light):
            return self.logo_light
        return None

    def _footer(self, slide, dark=False):
        col = RGBColor(0xBD, 0xBD, 0xBD) if dark else MUTED
        y = SLIDE_H - Inches(0.62)
        logo = self._logo_path_for(dark)
        if logo:
            # ØS: slide-nummer nederst til venstre, officielt logo nederst til højre
            nb = self._box(slide, MARGIN, y, Inches(3), Inches(0.4))
            _txt(nb.text_frame, "%02d" % self._n, font=FONT_MONO, size=11,
                 color=col, mono_label=True)
            lw = Inches(1.55)                       # bredde; højde efter aspect
            from PIL import Image as _PILImage
            try:
                iw, ih = _PILImage.open(logo).size
                lh = Emu(int(lw * ih / iw))
            except Exception:
                lh = Inches(0.42)
            slide.shapes.add_picture(logo, SLIDE_W - MARGIN - lw,
                                     y + Inches(0.04), width=lw, height=lh)
        else:
            left = self._box(slide, MARGIN, y, Inches(8), Inches(0.4))
            _txt(left.text_frame, self.footer, font=FONT_MONO, size=11, color=col, mono_label=True)
            right = self._box(slide, SLIDE_W - Inches(2.5), y, Inches(1.84), Inches(0.4))
            _txt(right.text_frame, "%02d" % self._n, font=FONT_MONO, size=11, color=col,
                 align=PP_ALIGN.RIGHT, mono_label=True)

    def _accent(self, slide, y=Inches(2.02), x=None):
        """Kort grøn accent-bjælke under en slide-titel (kun brand-tema).
        x flugter med titel-tekstens venstrekant = MARGIN + tekstboksens 0.1" indryk."""
        if not self.brand_accent:
            return
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     x if x is not None else MARGIN + Inches(0.1), y,
                                     Inches(0.7), Pt(6))
        _solid(bar, self.brand); _no_border(bar); bar.shadow.inherit = False

    def _kicker(self, slide, text, top, color=None):
        if not self.show_eyebrow:
            return None                      # temaet skjuler kicker-labels
        tb = self._box(slide, MARGIN, top, Inches(11), Inches(0.4))
        _txt(tb.text_frame, text, font=FONT_MONO, size=14, bold=True,
             color=color or self.kicker_color, mono_label=True)
        return tb

    # ============================================================ SLIDE-TYPER

    def title(self, title, lead=None, kicker=None):
        slide = self._new(CANVAS)
        # hvidt frame med stor shadow — teksten flyder inde i rammen
        fr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, Inches(1.7),
                                    Inches(11.0), Inches(4.0))
        _solid(fr, PAPER)
        _border(fr, INK, 3.0)
        fr.shadow.inherit = False
        # Brand-tema: grøn offset-shadow på titel-rammen (i stedet for grøn topkant)
        _offset_shadow(fr, self.brand if self.brand_accent else INK, dist_pt=8)
        blocks = []
        if kicker and self.show_eyebrow:
            blocks.append(dict(text=kicker, mono=True, size=14, bold=True,
                               color=self.kicker_color, after=18))
        blocks.append(dict(text=title, size=58, bold=True, color=INK, line=1.02, after=20))
        if lead:
            blocks.append(dict(text=lead, size=22, color=INK_SOFT, line=1.25))
        _fill(fr, blocks, anchor=MSO_ANCHOR.MIDDLE, ml=0.6, mr=0.6, mt=0.55, mb=0.55)
        self._footer(slide)
        return slide

    def section(self, title, num=None, variant=None):
        variant = variant or self.section_default
        bg = {"dark": INK, "brand": self.brand, "yellow": ACCENT}.get(variant, CANVAS)
        slide = self._new(bg)
        dark = variant in ("dark", "brand")
        txt_col = WHITE if dark else INK
        if num:
            nb = self._box(slide, MARGIN, Inches(2.3), Inches(4), Inches(0.6))
            _txt(nb.text_frame, num, font=FONT_MONO, size=26, bold=True, color=txt_col, mono_label=True)
        tb = self._box(slide, MARGIN, Inches(3.0), Inches(11.5), Inches(2.2))
        _txt(tb.text_frame, title, size=72, bold=True, color=txt_col, line_pct=1.0)
        self._footer(slide, dark=dark)
        return slide

    def content(self, title, bullets, kicker=None, accent_index=None):
        """bullets: liste af str. accent_index: index på bullet med rød markør."""
        slide = self._new(CANVAS)
        self._kicker(slide, kicker or "PRINCIP", Inches(0.66))
        t = self._box(slide, MARGIN, Inches(1.15), Inches(11.5), Inches(1.2))
        _txt(t.text_frame, title, size=42, bold=True, color=INK, line_pct=1.05)
        self._accent(slide)
        # bullets
        top = Inches(2.7)
        for i, b in enumerate(bullets):
            sq_color = self.signal if i == accent_index else INK
            sq = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, top + Pt(6),
                                        Pt(11), Pt(11))
            _solid(sq, sq_color)
            _no_border(sq)
            sq.shadow.inherit = False
            tb = self._box(slide, MARGIN + Inches(0.45), top, Inches(10.5), Inches(0.7))
            _txt(tb.text_frame, b, size=22, color=INK_SOFT, line_pct=1.15)
            top += Inches(0.72)
        self._footer(slide)
        return slide

    def split(self, title, left, right, kicker=None):
        """left/right: dict {'label': str, 'body': str}."""
        slide = self._new(CANVAS)
        self._kicker(slide, kicker or "OVERSIGT", Inches(0.66))
        t = self._box(slide, MARGIN, Inches(1.15), Inches(11.5), Inches(1.2))
        _txt(t.text_frame, title, size=42, bold=True, color=INK, line_pct=1.05)
        self._accent(slide)
        colw = Inches(5.3)
        # divider
        dv = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.6), Inches(3.0),
                                    Pt(1), Inches(3.3))
        _solid(dv, LINE_SOFT); _no_border(dv); dv.shadow.inherit = False
        for col, x in ((left, MARGIN), (right, Inches(7.0))):
            lb = self._box(slide, x, Inches(3.0), colw, Inches(0.4))
            _txt(lb.text_frame, col["label"], font=FONT_MONO, size=13, bold=True,
                 color=MUTED, mono_label=True)
            bd = self._box(slide, x, Inches(3.5), colw, Inches(2.7))
            _txt(bd.text_frame, col["body"], size=22, color=INK_SOFT, line_pct=1.3)
        self._footer(slide)
        return slide

    def statement(self, quote, source=None):
        slide = self._new(CANVAS)
        # signal-farvet venstre-kant
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, Inches(2.4),
                                     Pt(6), Inches(2.6))
        _solid(bar, self.signal); _no_border(bar); bar.shadow.inherit = False
        q = self._box(slide, MARGIN + Inches(0.4), Inches(2.4), Inches(10.5), Inches(2.6))
        q.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        _txt(q.text_frame, quote, size=44, bold=True, color=INK, line_pct=1.1)
        if source:
            s = self._box(slide, MARGIN + Inches(0.4), Inches(5.1), Inches(10), Inches(0.5))
            _txt(s.text_frame, source, font=FONT_MONO, size=15, bold=True,
                 color=MUTED, mono_label=True)
        self._footer(slide)
        return slide

    def stats(self, cards, title=None, kicker=None, primary_index=None):
        """cards: liste af dict {'num': str, 'label': str, 'trend': 'up'|'down'|None}.
        primary_index: index på det fremhævede kort (gul + rød shadow)."""
        slide = self._new(CANVAS)
        self._kicker(slide, kicker or "RESULTATER", Inches(0.66))
        if title:
            t = self._box(slide, MARGIN, Inches(1.15), Inches(11.5), Inches(1.0))
            _txt(t.text_frame, title, size=42, bold=True, color=INK)
            self._accent(slide)
        n = len(cards)
        gap = Inches(0.4)
        total_w = SLIDE_W - 2 * MARGIN
        cw = Emu(int((total_w - gap * (n - 1)) / n))
        top = Inches(3.0)
        ch = Inches(2.6)
        for i, c in enumerate(cards):
            x = Emu(int(MARGIN + i * (cw + gap)))
            primary = i == primary_index
            card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, top, cw, ch)
            _solid(card, ACCENT if primary else PAPER)
            _border(card, INK, 2.0)
            card.shadow.inherit = False
            _offset_shadow(card, self.signal if primary else INK, dist_pt=5)
            # tal
            num_col = INK
            if c.get("trend") == "up": num_col = self.success
            if c.get("trend") == "down": num_col = self.signal
            nb = self._box(slide, x + Inches(0.3), top + Inches(0.35),
                           Emu(int(cw - Inches(0.6))), Inches(1.3))
            # Auto-skrump: lange værdier ("1,2 mio") krymper til at passe i kortet
            # uanset skrifttype, i stedet for at ombryde/flyde over.
            nb.text_frame.word_wrap = False
            nb.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            _txt(nb.text_frame, c["num"], font=FONT_MONO, size=60, bold=True, color=num_col)
            lb = self._box(slide, x + Inches(0.3), top + Inches(1.75),
                           Emu(int(cw - Inches(0.6))), Inches(0.7))
            _txt(lb.text_frame, c["label"], font=FONT_MONO, size=13, bold=True,
                 color=MUTED, mono_label=True, line_pct=1.1)
        self._footer(slide)
        return slide

    def table(self, headers, rows, title=None, kicker=None,
              numeric_cols=(), highlight_row=None):
        """headers: list[str]. rows: list[list[str]]. numeric_cols: indices højrejusteret+mono."""
        slide = self._new(CANVAS)
        self._kicker(slide, kicker or "OVERSIGT", Inches(0.66))
        if title:
            t = self._box(slide, MARGIN, Inches(1.15), Inches(11.5), Inches(1.0))
            _txt(t.text_frame, title, size=42, bold=True, color=INK)
            self._accent(slide)
        ncols = len(headers)
        nrows = len(rows) + 1
        tbl_w = SLIDE_W - 2 * MARGIN
        gtbl = slide.shapes.add_table(nrows, ncols, MARGIN, Inches(2.7),
                                      tbl_w, Inches(0.6 * nrows)).table
        # disable PPT auto-styling
        tbl_el = gtbl._tbl
        for el in tbl_el.findall(qn('a:tblPr')):
            el.set('firstRow', '0'); el.set('bandRow', '0')
        # header
        for c, h in enumerate(headers):
            cell = gtbl.cell(0, c)
            _solid(cell, INK)
            cell.margin_left = Inches(0.18); cell.margin_top = Inches(0.08)
            _txt(cell.text_frame, h, font=FONT_MONO, size=12, bold=True, color=WHITE,
                 mono_label=True,
                 align=PP_ALIGN.RIGHT if c in numeric_cols else PP_ALIGN.LEFT)
        # body
        for r, row in enumerate(rows, start=1):
            hl = (r - 1) == highlight_row
            for c, val in enumerate(row):
                cell = gtbl.cell(r, c)
                _solid(cell, ACCENT if hl else PAPER)
                cell.margin_left = Inches(0.18); cell.margin_top = Inches(0.06)
                is_num = c in numeric_cols
                _txt(cell.text_frame, val,
                     font=FONT_MONO if is_num else FONT_SANS,
                     size=16, color=INK_SOFT,
                     align=PP_ALIGN.RIGHT if is_num else PP_ALIGN.LEFT)
        self._footer(slide)
        return slide

    def closing(self, title, cta=None, kicker=None):
        slide = self._new(CANVAS)
        self._kicker(slide, kicker or "TAK", Inches(2.0))
        t = self._box(slide, MARGIN, Inches(2.5), Inches(11.5), Inches(1.6))
        _txt(t.text_frame, title, size=66, bold=True, color=INK)
        if cta:
            btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, Inches(4.5),
                                         Inches(5), Inches(0.8))
            try:
                btn.adjustments[0] = 0.08
            except Exception:
                pass
            _solid(btn, self.cta_bg)
            _border(btn, self.cta_bg, 2.0)
            btn.shadow.inherit = False
            _offset_shadow(btn, self.signal, dist_pt=4)
            tf = btn.text_frame
            tf.word_wrap = True
            _txt(tf, cta, font=FONT_MONO, size=15, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, mono_label=True)
        self._footer(slide)
        return slide

    # ============================================================ PRÆSENTATIONS-ELEMENTER

    def _head(self, slide, kicker, title, kicker_default="OVERSIGT"):
        self._kicker(slide, kicker or kicker_default, Inches(0.66))
        if title:
            t = self._box(slide, MARGIN, Inches(1.15), Inches(11.5), Inches(1.0))
            _txt(t.text_frame, title, size=42, bold=True, color=INK, line_pct=1.05)
            self._accent(slide)

    def _img(self, slide, l, t, w, h, label="BILLEDE", path=None):
        """Billed-felt: ægte billede hvis path, ellers grå placeholder m. mono-label."""
        if path:
            slide.shapes.add_picture(path, l, t, w, h)
            box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
            _no_fill(box); _border(box, INK, 2.0); box.shadow.inherit = False
            return
        ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
        _solid(ph, LINE_SOFT); _border(ph, INK, 2.0); ph.shadow.inherit = False
        tf = ph.text_frame; tf.word_wrap = True
        _txt(tf, label, font=FONT_MONO, size=12, bold=True, color=MUTED,
             align=PP_ALIGN.CENTER, mono_label=True)

    def cover(self, title, lead=None, kicker=None, image=None):
        slide = self._new(CANVAS)
        self._img(slide, 0, 0, SLIDE_W, SLIDE_H, "BAGGRUNDSBILLEDE 16:9", image)
        # overlay-boks med hård shadow — teksten flyder inde i boksen
        ov = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, Inches(3.3),
                                    Inches(9.0), Inches(3.4))
        _solid(ov, PAPER); _border(ov, INK, 3.0); ov.shadow.inherit = False
        _offset_shadow(ov, INK, dist_pt=8)
        blocks = []
        if kicker and self.show_eyebrow:
            blocks.append(dict(text=kicker, mono=True, size=13, bold=True,
                               color=self.kicker_color, after=14))
        blocks.append(dict(text=title, size=44, bold=True, color=INK, line=1.04, after=14))
        if lead:
            blocks.append(dict(text=lead, size=18, color=INK_SOFT, line=1.2))
        _fill(ov, blocks, anchor=MSO_ANCHOR.MIDDLE, ml=0.55, mr=0.55, mt=0.5, mb=0.5)
        return slide

    def image_full(self, caption=None, kicker=None, image=None):
        slide = self._new(CANVAS)
        self._img(slide, 0, 0, SLIDE_W, SLIDE_H, "FULL-BLEED BILLEDE", image)
        if caption:
            strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.2),
                                           Inches(9), Inches(0.7))
            _solid(strip, INK); _no_border(strip); strip.shadow.inherit = False
            cap = self._box(slide, Inches(0.3), Inches(6.32), Inches(8.5), Inches(0.5))
            p = cap.text_frame.paragraphs[0]
            if kicker:
                r1 = p.add_run(); r1.text = kicker.upper() + "   "
                r1.font.name = FONT_MONO; r1.font.size = Pt(12); r1.font.bold = True
                r1.font.color.rgb = self.signal; _letter_spacing(r1, 1.0)
            r2 = p.add_run(); r2.text = caption
            r2.font.name = FONT_SANS; r2.font.size = Pt(18); r2.font.color.rgb = WHITE
        return slide

    def image_text(self, title, bullets, kicker=None, image=None,
                   image_left=False, accent_index=None):
        slide = self._new(CANVAS)
        self._head(slide, kicker, title, "CASE")
        colw = Inches(5.5); top = Inches(2.5); ch = Inches(4.2)
        img_x = MARGIN if image_left else Inches(7.13)
        txt_x = Inches(7.13) if image_left else MARGIN
        self._img(slide, img_x, top, colw, ch, "BILLEDE", image)
        by = top
        for i, b in enumerate(bullets):
            sq = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, txt_x, by + Pt(6), Pt(11), Pt(11))
            _solid(sq, self.signal if i == accent_index else INK)
            _no_border(sq); sq.shadow.inherit = False
            tb = self._box(slide, txt_x + Inches(0.4), by, Emu(int(colw - Inches(0.4))), Inches(0.9))
            _txt(tb.text_frame, b, size=20, color=INK_SOFT, line_pct=1.15)
            by += Inches(0.95)
        self._footer(slide)
        return slide

    def image_grid(self, title=None, n=6, cols=3, kicker=None, images=None):
        slide = self._new(CANVAS)
        self._head(slide, kicker, title, "GALLERI")
        images = images or [None] * n
        gap = Inches(0.25); top = Inches(2.5)
        rows = (n + cols - 1) // cols
        gw = SLIDE_W - 2 * MARGIN
        cw = Emu(int((gw - gap * (cols - 1)) / cols))
        gh = Inches(4.2)
        chh = Emu(int((gh - gap * (rows - 1)) / rows))
        for i in range(n):
            r, c = divmod(i, cols)
            x = Emu(int(MARGIN + c * (cw + gap)))
            y = Emu(int(top + r * (chh + gap)))
            self._img(slide, x, y, cw, chh, "%02d" % (i + 1), images[i] if i < len(images) else None)
        self._footer(slide)
        return slide

    def agenda(self, items, title="Agenda", kicker=None, current=None):
        slide = self._new(CANVAS)
        self._head(slide, kicker, title, "INDHOLD")
        y = Inches(2.4)
        for i, it in enumerate(items):
            is_cur = i == current
            nb = self._box(slide, MARGIN, y, Inches(1), Inches(0.6))
            _txt(nb.text_frame, "%02d" % (i + 1), font=FONT_MONO, size=22, bold=True,
                 color=self.signal if is_cur else MUTED, mono_label=True)
            tb = self._box(slide, MARGIN + Inches(1.1), y - Inches(0.02), Inches(10), Inches(0.6))
            _txt(tb.text_frame, it, size=26, bold=is_cur, color=INK)
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, y + Inches(0.7),
                                          SLIDE_W - 2 * MARGIN, Pt(1))
            _solid(line, LINE_SOFT); _no_border(line); line.shadow.inherit = False
            y += Inches(0.95)
        self._footer(slide)
        return slide

    def cards(self, items, title=None, kicker=None, accent_index=None):
        """items: liste af dict {'title': str, 'body': str}."""
        slide = self._new(CANVAS)
        self._head(slide, kicker, title, "FUNKTIONER")
        n = len(items); gap = Inches(0.35); top = Inches(2.6)
        gw = SLIDE_W - 2 * MARGIN
        cw = Emu(int((gw - gap * (n - 1)) / n)); ch = Inches(3.6)
        for i, it in enumerate(items):
            acc = i == accent_index
            x = Emu(int(MARGIN + i * (cw + gap)))
            card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, top, cw, ch)
            _solid(card, ACCENT if acc else PAPER); _border(card, INK, 2.0)
            card.shadow.inherit = False; _offset_shadow(card, INK, dist_pt=3)
            ic = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.3), top + Inches(0.3),
                                        Inches(0.5), Inches(0.5))
            _solid(ic, self.signal if acc else INK); _no_border(ic); ic.shadow.inherit = False
            tt = self._box(slide, x + Inches(0.3), top + Inches(1.0), Emu(int(cw - Inches(0.6))), Inches(0.6))
            _txt(tt.text_frame, it["title"], size=22, bold=True, color=INK)
            bd = self._box(slide, x + Inches(0.3), top + Inches(1.6), Emu(int(cw - Inches(0.6))), Inches(1.8))
            _txt(bd.text_frame, it["body"], size=16, color=INK_SOFT, line_pct=1.3)
        self._footer(slide)
        return slide

    def timeline(self, steps, title=None, kicker=None, current=None):
        """steps: liste af dict {'when': str, 'what': str, 'desc': str}."""
        slide = self._new(CANVAS)
        self._head(slide, kicker, title, "PROCES")
        n = len(steps); top = Inches(3.0)
        gw = SLIDE_W - 2 * MARGIN
        step_w = Emu(int(gw / n))
        for i, s in enumerate(steps):
            x = Emu(int(MARGIN + i * step_w))
            cur = i == current
            # connector
            if i < n - 1:
                ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(int(x + Inches(0.24))),
                                            top + Inches(0.1), Emu(int(step_w - Inches(0.24))), Pt(3))
                _solid(ln, INK); _no_border(ln); ln.shadow.inherit = False
            dot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, top, Inches(0.26), Inches(0.26))
            _solid(dot, self.signal if cur else PAPER)
            _border(dot, self.signal if cur else INK, 3.0); dot.shadow.inherit = False
            wb = self._box(slide, x, top + Inches(0.45), Emu(int(step_w - Inches(0.2))), Inches(0.4))
            _txt(wb.text_frame, s["when"], font=FONT_MONO, size=13, bold=True, color=MUTED, mono_label=True)
            tb = self._box(slide, x, top + Inches(0.85), Emu(int(step_w - Inches(0.2))), Inches(0.5))
            _txt(tb.text_frame, s["what"], size=18, bold=True, color=INK)
            db = self._box(slide, x, top + Inches(1.4), Emu(int(step_w - Inches(0.25))), Inches(1.6))
            _txt(db.text_frame, s.get("desc", ""), size=14, color=INK_SOFT, line_pct=1.25)
        self._footer(slide)
        return slide

    def comparison(self, left, right, title=None, kicker=None, win="right"):
        """left/right: dict {'head': str, 'items': [(bool_yes, str), ...]}."""
        slide = self._new(CANVAS)
        self._head(slide, kicker, title, "VALG")
        top = Inches(2.6); ch = Inches(4.0)
        cw = Inches(5.8); gap = Inches(0.4)
        for side, col in (("left", left), ("right", right)):
            x = MARGIN if side == "left" else Emu(int(MARGIN + cw + gap))
            card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, top, cw, ch)
            _solid(card, PAPER); _border(card, INK, 2.0); card.shadow.inherit = False
            if side == win:
                _offset_shadow(card, self.signal, dist_pt=5)
            hb = self._box(slide, x + Inches(0.3), top + Inches(0.25), Emu(int(cw - Inches(0.6))), Inches(0.5))
            _txt(hb.text_frame, col["head"], font=FONT_MONO, size=14, bold=True, color=INK, mono_label=True)
            iy = top + Inches(1.0)
            for yes, text in col["items"]:
                mk = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.3), iy + Pt(3),
                                            Inches(0.22), Inches(0.22))
                _solid(mk, self.success if yes else self.signal); _no_border(mk); mk.shadow.inherit = False
                tb = self._box(slide, x + Inches(0.7), iy, Emu(int(cw - Inches(1.0))), Inches(0.5))
                _txt(tb.text_frame, text, size=17, color=INK_SOFT)
                iy += Inches(0.6)
        self._footer(slide)
        return slide

    def bignum(self, value, sub=None, kicker=None, trend=None):
        slide = self._new(CANVAS)
        self._kicker(slide, kicker or "NØGLETAL", Inches(1.6))
        col = INK
        if trend == "up": col = self.success
        if trend == "down": col = self.signal
        nb = self._box(slide, MARGIN, Inches(2.1), Inches(11.5), Inches(2.8))
        nb.text_frame.word_wrap = False
        nb.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE   # langt tal krymper til at passe
        _txt(nb.text_frame, value, font=FONT_MONO, size=150, bold=True, color=col, line_pct=0.9)
        if sub:
            sb = self._box(slide, MARGIN, Inches(5.2), Inches(9), Inches(1.0))
            _txt(sb.text_frame, sub, size=24, color=INK_SOFT, line_pct=1.2)
        self._footer(slide)
        return slide

    def bars(self, rows, title=None, kicker=None, accent_index=None):
        """rows: liste af dict {'label': str, 'pct': 0-100, 'value': str}."""
        slide = self._new(CANVAS)
        self._head(slide, kicker, title, "FORDELING")
        top = Inches(2.7); rh = Inches(0.55); gap = Inches(0.35)
        lbl_w = Inches(2.2); val_w = Inches(1.1)
        track_x = MARGIN + lbl_w
        track_w = SLIDE_W - MARGIN - track_x - val_w - Inches(0.2)
        for i, r in enumerate(rows):
            y = Emu(int(top + i * (rh + gap)))
            lb = self._box(slide, MARGIN, y + Inches(0.08), lbl_w, Inches(0.4))
            _txt(lb.text_frame, r["label"], size=17, color=INK)
            track = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, track_x, y, track_w, rh)
            _solid(track, PAPER); _border(track, INK, 2.0); track.shadow.inherit = False
            fill_w = Emu(int(track_w * max(0, min(100, r["pct"])) / 100))
            if fill_w > 0:
                fill = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, track_x, y, fill_w, rh)
                _solid(fill, self.signal if i == accent_index else INK)
                _no_border(fill); fill.shadow.inherit = False
            vb = self._box(slide, Emu(int(track_x + track_w + Inches(0.15))), y + Inches(0.08), val_w, Inches(0.4))
            _txt(vb.text_frame, r["value"], font=FONT_MONO, size=17, bold=True, color=INK,
                 align=PP_ALIGN.RIGHT)
        self._footer(slide)
        return slide

    def testimonial(self, quote, name, role, kicker=None, image=None):
        slide = self._new(CANVAS)
        self._kicker(slide, kicker or "UDTALELSE", Inches(0.66))
        self._img(slide, MARGIN, Inches(2.0), Inches(2.6), Inches(2.6), "PORTRÆT", image)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.0), Inches(2.0), Pt(6), Inches(2.0))
        _solid(bar, self.signal); _no_border(bar); bar.shadow.inherit = False
        q = self._box(slide, Inches(4.3), Inches(2.0), Inches(8.0), Inches(2.2))
        _txt(q.text_frame, quote, size=30, bold=True, color=INK, line_pct=1.2)
        nb = self._box(slide, Inches(4.3), Inches(4.3), Inches(8), Inches(0.5))
        _txt(nb.text_frame, name, size=18, bold=True, color=INK)
        rb = self._box(slide, Inches(4.3), Inches(4.8), Inches(8), Inches(0.4))
        _txt(rb.text_frame, role, font=FONT_MONO, size=13, bold=True, color=MUTED, mono_label=True)
        self._footer(slide)
        return slide

    def save(self, path):
        self.prs.save(path)
        return path
